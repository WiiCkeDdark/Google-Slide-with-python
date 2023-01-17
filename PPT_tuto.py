from pptx import Presentation
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.chart import XL_LEGEND_POSITION

# create function to create first slide


def create_first_slide(prs, title, subtitle):
    first_slide = prs.slides.add_slide(prs.slide_layouts[0])
    first_slide.shapes.title.text = title
    first_slide.placeholders[1].text = subtitle
    return first_slide

# create function to create new slide


def create_new_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    return slide

# create function to create chart


def create_chart(slide, chart_type, x, y, cx, cy, chart_data):
    graphic_frame = slide.shapes.add_chart(
        chart_type, x, y, cx, cy, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return chart


if __name__ == '__main__':
    # create PPT
    prs = Presentation()
    create_first_slide(prs, "Tuto PPT Python",
                       "Création de PowerPoint avec Python - by Antoine")
    slide = create_new_slide(prs, "Graphique")
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
    chart_data = ChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Q1 Sales', (19.2, 21.4, 16.7))
    chart_data.add_series('Q2 Sales', (22.3, 28.6, 15.2))
    chart_data.add_series('Q3 Sales', (20.4, 26.3, 14.2))
    chart = create_chart(
        slide, XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    # create last slide
    slide = create_new_slide(prs, "Fin du PPT")
    prs.save('PPT_tuto.pptx')
